"""
Blur sensitive UI regions (names, NIDA, phone, TIN) in portal screenshots,
then swap those images into the SEM2 proposal.docx without changing fonts/layout.
Optionally embed a few blurred shots into Ushirika_Group15.pptx.
"""

from __future__ import annotations

import io
import shutil
import zipfile
from pathlib import Path

from PIL import Image, ImageFilter

ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "Presentation" / "screenshots"
BLURRED = SHOTS / "blurred"

PROPOSAL = Path(r"C:\Users\USER\OneDrive\Desktop\SEM 2\Data Science Project\Group 15 proposal.docx")
PPTX = Path(r"C:\Users\USER\OneDrive\Desktop\SEM 2\Data Science Project\Ushirika_Group15.pptx")

# Relative boxes as (left, top, right, bottom) fractions of image size.
# Tuned for Ushirika portal layouts to cover names / NIDA / phone / TIN.
BLUR_REGIONS: dict[str, list[tuple[float, float, float, float]]] = {
    "01_login.png": [
        (0.28, 0.42, 0.72, 0.52),  # ID field value
    ],
    "02_sme_dashboard.png": [
        (0.12, 0.08, 0.42, 0.16),  # topbar name
        (0.55, 0.22, 0.92, 0.55),  # score/components personal notes
    ],
    "03_sme_transactions.png": [
        (0.12, 0.08, 0.42, 0.16),
        (0.08, 0.38, 0.95, 0.88),  # table body with names/TIN
    ],
    "04_sme_upload.png": [
        (0.12, 0.08, 0.42, 0.16),
    ],
    "05_lender_portfolio.png": [
        (0.12, 0.08, 0.45, 0.16),
        (0.08, 0.32, 0.45, 0.90),  # SME list names + NIDA
    ],
    "06_lender_sme_detail.png": [
        (0.12, 0.08, 0.45, 0.16),
        (0.08, 0.28, 0.95, 0.48),  # profile NIDA/name/phone/email
        (0.08, 0.32, 0.40, 0.90),
    ],
    "07_admin_accounts.png": [
        (0.12, 0.08, 0.45, 0.16),
        (0.08, 0.36, 0.95, 0.90),  # accounts table
    ],
    "08_admin_create.png": [
        (0.12, 0.08, 0.45, 0.16),
        (0.28, 0.30, 0.72, 0.78),  # form fields with identity
    ],
    "09_lender_ml_metrics.png": [
        (0.12, 0.08, 0.45, 0.16),
        (0.08, 0.22, 0.95, 0.40),
    ],
    "10_lender_signals.png": [
        (0.12, 0.08, 0.45, 0.16),
        (0.08, 0.22, 0.95, 0.40),
    ],
    "10b_lender_signals_full.png": [
        (0.12, 0.08, 0.45, 0.16),
        (0.08, 0.22, 0.95, 0.40),
    ],
    "11_lender_recent_tx.png": [
        (0.12, 0.08, 0.45, 0.16),
        (0.08, 0.35, 0.95, 0.90),
    ],
}

# proposal.docx word/media file → screenshot (skip cover/logo if unknown)
# Inspected previously: image2.. often figures. We map by replacing content of
# matching figure screenshots while preserving zip entry names.
PROPOSAL_MEDIA_MAP = {
    # Keep flexible: replace by scanning existing media and matching size/order later
}


def blur_image(src: Path, regions: list[tuple[float, float, float, float]]) -> Image.Image:
    img = Image.open(src).convert("RGB")
    w, h = img.size
    out = img.copy()
    for l, t, r, b in regions:
        box = (
            max(0, int(l * w)),
            max(0, int(t * h)),
            min(w, int(r * w)),
            min(h, int(b * h)),
        )
        if box[2] <= box[0] or box[3] <= box[1]:
            continue
        crop = img.crop(box).filter(ImageFilter.GaussianBlur(radius=22))
        # Soft cover so digits/names stay unreadable
        cover = Image.new("RGB", crop.size, (225, 225, 225))
        crop = Image.blend(crop, cover, alpha=0.35)
        out.paste(crop, box[:2])
    return out


def build_blurred_shots() -> dict[str, Path]:
    BLURRED.mkdir(parents=True, exist_ok=True)
    outputs: dict[str, Path] = {}
    for name, regions in BLUR_REGIONS.items():
        src = SHOTS / name
        if not src.exists():
            print(f"skip missing {name}")
            continue
        dest = BLURRED / name
        blur_image(src, regions).save(dest, format="PNG", optimize=True)
        outputs[name] = dest
        print(f"blurred {name} -> {dest}")
    return outputs


def _png_bytes(path: Path) -> bytes:
    img = Image.open(path).convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=True)
    return buf.getvalue()


def replace_proposal_media(blurred: dict[str, Path]) -> None:
    if not PROPOSAL.exists():
        raise FileNotFoundError(PROPOSAL)

    backup = PROPOSAL.with_suffix(".docx.bak")
    if not backup.exists():
        shutil.copy2(PROPOSAL, backup)
        print(f"backup {backup}")

    # Map figure-oriented screenshots onto media parts by size ranking of existing PNGs
    # (largest portal screenshots replace the largest figure media, excluding tiny logos).
    with zipfile.ZipFile(PROPOSAL, "r") as zf:
        media = [n for n in zf.namelist() if n.startswith("word/media/") and n.lower().endswith((".png", ".jpeg", ".jpg"))]
        sizes = []
        for name in media:
            data = zf.read(name)
            try:
                im = Image.open(io.BytesIO(data))
                sizes.append((im.size[0] * im.size[1], name, data))
            except Exception:
                continue
        sizes.sort(reverse=True)

    # Prefer known portal shots in a stable figure order
    preferred = [
        "02_sme_dashboard.png",
        "01_login.png",
        "03_sme_transactions.png",
        "04_sme_upload.png",
        "05_lender_portfolio.png",
        "06_lender_sme_detail.png",
        "07_admin_accounts.png",
        "08_admin_create.png",
        "09_lender_ml_metrics.png",
        "10_lender_signals.png",
        "11_lender_recent_tx.png",
        "10b_lender_signals_full.png",
    ]
    shot_paths = [blurred[n] for n in preferred if n in blurred]
    if not shot_paths:
        raise RuntimeError("No blurred screenshots available")

    # Replace the largest N media entries (skip the absolute largest if it looks like a cover)
    # Keep first (cover) untouched when more than 8 media exist.
    candidates = sizes[1:] if len(sizes) > 8 else sizes
    replacements = {}
    for i, shot in enumerate(shot_paths):
        if i >= len(candidates):
            break
        _, media_name, _ = candidates[i]
        replacements[media_name] = _png_bytes(shot)
        print(f"proposal {media_name} <- {shot.name}")

    tmp = PROPOSAL.with_suffix(".docx.tmp")
    with zipfile.ZipFile(PROPOSAL, "r") as zin, zipfile.ZipFile(tmp, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = replacements.get(item.filename, zin.read(item.filename))
            zout.writestr(item, data)
    tmp.replace(PROPOSAL)
    print(f"updated {PROPOSAL}")


def update_pptx_images(blurred: dict[str, Path]) -> None:
    """Add blurred screenshots to the SEM2 PPT without rewriting text/fonts."""
    if not PPTX.exists():
        print(f"skip pptx missing {PPTX}")
        return
    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
    except ImportError:
        print("python-pptx not installed; skip PPT image embed")
        return

    backup = PPTX.with_suffix(".pptx.bak")
    if not backup.exists():
        shutil.copy2(PPTX, backup)

    prs = Presentation(str(PPTX))
    # If slides already have pictures, replace picture blobs in zip instead.
    with zipfile.ZipFile(PPTX, "r") as zf:
        media = [n for n in zf.namelist() if n.startswith("ppt/media/")]
    if media:
        preferred = [
            "02_sme_dashboard.png",
            "05_lender_portfolio.png",
            "07_admin_accounts.png",
            "03_sme_transactions.png",
            "01_login.png",
            "08_admin_create.png",
        ]
        shots = [blurred[n] for n in preferred if n in blurred]
        replacements = {}
        for i, m in enumerate(sorted(media)):
            if i >= len(shots):
                break
            replacements[m] = _png_bytes(shots[i])
            print(f"pptx {m} <- {shots[i].name}")
        tmp = PPTX.with_suffix(".pptx.tmp")
        with zipfile.ZipFile(PPTX, "r") as zin, zipfile.ZipFile(tmp, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = replacements.get(item.filename, zin.read(item.filename))
                zout.writestr(item, data)
        tmp.replace(PPTX)
        print(f"updated media in {PPTX}")
        return

    # No media yet: place one picture on each of the last few content slides
    # without changing existing text boxes (append picture shapes only).
    preferred = [
        "01_login.png",
        "02_sme_dashboard.png",
        "05_lender_portfolio.png",
        "07_admin_accounts.png",
    ]
    shots = [blurred[n] for n in preferred if n in blurred]
    slide_indexes = list(range(max(0, len(prs.slides) - len(shots)), len(prs.slides)))
    for slide_idx, shot in zip(slide_indexes, shots):
        slide = prs.slides[slide_idx]
        # Place lower-right so titles/body text stay visible
        slide.shapes.add_picture(str(shot), Inches(5.4), Inches(3.4), width=Inches(4.2))
        print(f"pptx slide {slide_idx + 1} + {shot.name}")
    prs.save(str(PPTX))
    print(f"saved {PPTX}")


def main() -> None:
    blurred = build_blurred_shots()
    replace_proposal_media(blurred)
    update_pptx_images(blurred)


if __name__ == "__main__":
    main()
