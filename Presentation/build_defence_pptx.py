"""
Ushirika / Group 15 defence presentation — 9 slides.
Readable on a projector, human wording, clean layout.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Ushirika_Group15_Defence_Presentation.pptx"
OUT_ALT = Path(__file__).resolve().parent / "Ushirika_Group15_Defence_v13.pptx"
OUT_SEM2 = Path(r"C:\Users\USER\OneDrive\Desktop\SEM 2\Data Science Project\Ushirika_Group15.pptx")

FOREST = RGBColor(0x0B, 0x3D, 0x2E)
FOREST_DEEP = RGBColor(0x06, 0x28, 0x20)
LAGOON = RGBColor(0x1A, 0x7A, 0x6D)
LAGOON_BRIGHT = RGBColor(0x23, 0x96, 0x88)
MIST = RGBColor(0xE6, 0xEE, 0xF0)
PAPER = RGBColor(0xF0, 0xF5, 0xF6)
INK = RGBColor(0x0F, 0x28, 0x30)
MUTED = RGBColor(0x3D, 0x52, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x1A, 0x6B, 0x45)
SOFT = RGBColor(0xD8, 0xE8, 0xE4)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 9

MEMBERS = [
    "Herman Edward Mkumbwa",
    "Raymond Elphance Tungaraza",
    "Edwin Celestin Silayo",
    "Grace Joachim Mohammed",
    "Priscila Nestor Mpembela",
]


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill(s, color)
    return s


def round_rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill(s, color)
    try:
        s.adjustments[0] = 0.08
    except Exception:
        pass
    return s


def set_run(run, text, size=20, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def textbox(slide, left, top, width, height, text, size=20, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n") if text is not None else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        set_run(run, line, size, bold, color, font)
    return box


def bullets(slide, left, top, width, height, items, size=18, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5)
        p.space_after = Pt(5)
        run = p.add_run()
        set_run(run, f"•  {item}", size, False, color)
    return box


def footer(slide, page):
    rect(slide, Inches(0), Inches(7.12), SLIDE_W, Inches(0.38), FOREST_DEEP)
    textbox(
        slide,
        Inches(0.45),
        Inches(7.16),
        Inches(10),
        Inches(0.3),
        "USHIRIKA  ·  Group 15  ·  EASTC",
        size=14,
        color=RGBColor(0xB8, 0xD0, 0xC8),
    )
    textbox(
        slide,
        Inches(11.2),
        Inches(7.16),
        Inches(1.7),
        Inches(0.3),
        f"{page}  /  {TOTAL}",
        size=14,
        color=RGBColor(0xB8, 0xD0, 0xC8),
        align=PP_ALIGN.RIGHT,
    )


def bg(slide, color=PAPER):
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, color)


def section_header(slide, kicker, title, subtitle=None):
    textbox(slide, Inches(0.55), Inches(0.18), Inches(12), Inches(0.32), kicker, size=16, bold=True, color=LAGOON)
    rect(slide, Inches(0.55), Inches(0.52), Inches(0.12), Inches(0.55), LAGOON_BRIGHT)
    textbox(
        slide,
        Inches(0.85),
        Inches(0.45),
        Inches(11.7),
        Inches(0.58),
        title,
        size=30,
        bold=True,
        color=FOREST,
        font="Georgia",
    )
    if subtitle:
        textbox(slide, Inches(0.55), Inches(1.08), Inches(12.2), Inches(0.38), subtitle, size=17, color=MUTED)


def metric(slide, left, top, w, h, value, label, accent=LAGOON):
    round_rect(slide, left, top, w, h, WHITE)
    rect(slide, left, top, Inches(0.12), h, accent)
    textbox(slide, left + Inches(0.28), top + Inches(0.22), w - Inches(0.4), Inches(0.45), value, size=28, bold=True, color=FOREST)
    textbox(slide, left + Inches(0.28), top + Inches(0.75), w - Inches(0.4), Inches(0.4), label, size=16, color=MUTED)


def met_badge(slide, left, top):
    round_rect(slide, left, top, Inches(0.95), Inches(0.36), SUCCESS)
    textbox(slide, left, top + Inches(0.02), Inches(0.95), Inches(0.32), "MET", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # —— 1 Title ——
    s = prs.slides.add_slide(blank)
    bg(s, FOREST_DEEP)
    rect(s, Inches(0), Inches(0), Inches(0.4), SLIDE_H, LAGOON_BRIGHT)
    rect(s, Inches(0.4), Inches(0), Inches(0.08), SLIDE_H, LAGOON)

    textbox(
        s,
        Inches(0.85),
        Inches(0.28),
        Inches(11.8),
        Inches(0.38),
        "EASTERN AFRICA STATISTICAL TRAINING CENTRE (EASTC)",
        size=18,
        bold=True,
        color=LAGOON_BRIGHT,
    )
    textbox(
        s,
        Inches(0.85),
        Inches(0.68),
        Inches(11.8),
        Inches(0.32),
        "Capstone Project Defence  ·  Bachelor of Data Science (Year III)  ·  2025/2026",
        size=16,
        color=MIST,
    )

    textbox(
        s,
        Inches(0.85),
        Inches(1.15),
        Inches(11.8),
        Inches(1.35),
        "Development of a Machine Learning-Based\n"
        "Credit Risk Assessment Model for SME\n"
        "Value Chain Financing",
        size=32,
        bold=True,
        color=WHITE,
        font="Georgia",
    )
    textbox(
        s,
        Inches(0.85),
        Inches(2.6),
        Inches(11.8),
        Inches(0.38),
        "A Case Study of Tanzania’s Supply Chains  ·  Platform: Ushirika",
        size=18,
        color=SOFT,
    )

    # Left: group + members
    round_rect(s, Inches(0.85), Inches(3.15), Inches(7.4), Inches(2.9), RGBColor(0x0A, 0x35, 0x2A))
    textbox(s, Inches(1.1), Inches(3.3), Inches(6.9), Inches(0.35), "GROUP 15  ·  Presented by", size=16, bold=True, color=LAGOON_BRIGHT)
    left_names = "\n".join(MEMBERS[:3])
    right_names = "\n".join(MEMBERS[3:])
    textbox(s, Inches(1.1), Inches(3.8), Inches(3.5), Inches(1.9), left_names, size=18, color=WHITE)
    textbox(s, Inches(4.7), Inches(3.8), Inches(3.3), Inches(1.9), right_names, size=18, color=WHITE)

    # Right: supervisor + link
    round_rect(s, Inches(8.5), Inches(3.15), Inches(4.3), Inches(2.9), RGBColor(0x0A, 0x35, 0x2A))
    textbox(s, Inches(8.75), Inches(3.35), Inches(3.9), Inches(0.3), "SUPERVISOR", size=15, bold=True, color=LAGOON_BRIGHT)
    textbox(s, Inches(8.75), Inches(3.75), Inches(3.9), Inches(0.55), "Mr. Rajabu Msangi", size=20, bold=True, color=WHITE)
    textbox(s, Inches(8.75), Inches(4.5), Inches(3.9), Inches(0.3), "LIVE PROJECT LINK", size=15, bold=True, color=LAGOON_BRIGHT)
    textbox(
        s,
        Inches(8.75),
        Inches(4.9),
        Inches(3.9),
        Inches(0.9),
        "ushirika-sme-portal\n.vercel.app",
        size=18,
        bold=True,
        color=SOFT,
    )

    textbox(s, Inches(0.85), Inches(6.65), Inches(10), Inches(0.35), "https://ushirika-sme-portal.vercel.app", size=16, color=RGBColor(0x9B, 0xC4, 0xBA))
    textbox(s, Inches(11.2), Inches(6.65), Inches(1.7), Inches(0.35), "1  /  9", size=16, color=RGBColor(0xB8, 0xD0, 0xC8), align=PP_ALIGN.RIGHT)

    # —— 2 Background ——
    s = prs.slides.add_slide(blank)
    bg(s)
    section_header(s, "BACKGROUND", "Why many Tanzanian SMEs still struggle to get fair credit")
    cards = [
        ("01", "Collateral-heavy lending", "Banks still ask for collateral and formal statements that most small traders simply do not have."),
        ("02", "Static risk ratios", "Old financial ratios often miss healthy businesses that work in informal or semi-formal supply chains."),
        ("03", "Unused payment signals", "How buyers and suppliers actually pay is useful risk data — but few systems capture and use it."),
        ("04", "The missing middle", "Without another way to score them, good SMEs stay locked out of formal finance."),
    ]
    for i, (num, title, body) in enumerate(cards):
        left = Inches(0.45 + (i % 2) * 6.4)
        top = Inches(1.65 + (i // 2) * 2.5)
        round_rect(s, left, top, Inches(6.15), Inches(2.25), WHITE)
        rect(s, left, top, Inches(0.14), Inches(2.25), LAGOON if i % 2 == 0 else FOREST)
        textbox(s, left + Inches(0.4), top + Inches(0.28), Inches(1.1), Inches(0.4), num, size=24, bold=True, color=LAGOON)
        textbox(s, left + Inches(1.4), top + Inches(0.32), Inches(4.4), Inches(0.4), title, size=20, bold=True, color=FOREST)
        textbox(s, left + Inches(0.4), top + Inches(0.9), Inches(5.4), Inches(1.15), body, size=18, color=MUTED)
    footer(s, 2)

    # —— 3 Aim ——
    s = prs.slides.add_slide(blank)
    bg(s)
    section_header(s, "AIM OF THE PROJECT", "What we set out to build")
    round_rect(s, Inches(0.45), Inches(1.6), Inches(12.4), Inches(1.45), WHITE)
    rect(s, Inches(0.45), Inches(1.6), Inches(0.14), Inches(1.45), LAGOON)
    textbox(s, Inches(0.8), Inches(1.75), Inches(11.8), Inches(0.3), "GENERAL OBJECTIVE", size=15, bold=True, color=LAGOON)
    textbox(
        s,
        Inches(0.8),
        Inches(2.15),
        Inches(11.8),
        Inches(0.7),
        "Build a working platform that reads supply-chain transactions and uses machine learning "
        "to give fairer, clearer credit risk scores for Tanzanian SMEs.",
        size=19,
        color=INK,
    )
    objs = [
        ("01", "Features from data", "Turn raw supply-chain transactions into useful predictors."),
        ("02", "Compare two models", "Test Random Forest against Logistic Regression."),
        ("03", "Prove with metrics", "Judge quality with Accuracy, Precision, Recall, F1 and ROC-AUC."),
    ]
    for i, (num, title, body) in enumerate(objs):
        left = Inches(0.45 + i * 4.2)
        round_rect(s, left, Inches(3.35), Inches(4.0), Inches(3.15), WHITE)
        rect(s, left, Inches(3.35), Inches(4.0), Inches(0.7), FOREST if i != 1 else LAGOON)
        textbox(s, left + Inches(0.15), Inches(3.5), Inches(3.7), Inches(0.45), f"{num}  {title}", size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(s, left + Inches(0.3), Inches(4.3), Inches(3.4), Inches(1.9), body, size=18, color=INK)
    footer(s, 3)

    # —— 4 What was built ——
    s = prs.slides.add_slide(blank)
    bg(s)
    section_header(s, "WHAT WAS BUILT", "Ushirika — a live credit platform for SMEs and lenders")
    layers = [
        ("Frontend", "Vite portal\nSME · Lender · Admin\nEnglish / Kiswahili", LAGOON),
        ("API", "FastAPI + JWT\nSecure login flows\nRole-based access", FOREST),
        ("ML Core", "Feature engineering\nRF (main) + LR\nPlain repayment signals", LAGOON_BRIGHT),
        ("Data & care", "SQL storage\nPII protected\nCareful loan caps", FOREST_DEEP),
    ]
    for i, (t, d, c) in enumerate(layers):
        left = Inches(0.4 + i * 3.25)
        round_rect(s, left, Inches(1.65), Inches(3.05), Inches(3.35), WHITE)
        rect(s, left, Inches(1.65), Inches(3.05), Inches(0.65), c)
        textbox(s, left + Inches(0.1), Inches(1.78), Inches(2.85), Inches(0.4), t, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(s, left + Inches(0.2), Inches(2.55), Inches(2.65), Inches(2.2), d, size=18, color=INK, align=PP_ALIGN.CENTER)
        if i < 3:
            textbox(s, left + Inches(2.9), Inches(3.0), Inches(0.4), Inches(0.4), "→", size=24, bold=True, color=LAGOON, align=PP_ALIGN.CENTER)
    round_rect(s, Inches(0.4), Inches(5.2), Inches(12.5), Inches(1.55), WHITE)
    textbox(
        s,
        Inches(0.7),
        Inches(5.4),
        Inches(12.0),
        Inches(1.2),
        "Lenders can search an SME, open the profile, and see score, risk band, plain signals and history.\n"
        "Registration already checks NIDA against date of birth, and uses region and district lists.\n"
        "The same portal works on desktop and on phone.",
        size=18,
        color=MUTED,
    )
    footer(s, 4)

    # —— 5 Literature ——
    s = prs.slides.add_slide(blank)
    bg(s)
    section_header(s, "LITERATURE REVIEW", "What prior work shows — and what is still missing here")
    cites = [
        ("Breiman (2001)", "Random Forests handle non-linear patterns better than a single model."),
        ("Lessmann et al. (2015)", "Ensemble classifiers beat classical baselines in credit scoring."),
        ("Khandani et al. (2010)", "Transaction data can stand in for creditworthiness when statements are thin."),
        ("Klapper (2006)", "Value-chain finance leans on buyer–supplier links, not only collateral."),
    ]
    for i, (a, b) in enumerate(cites):
        top = Inches(1.55 + i * 1.05)
        round_rect(s, Inches(0.45), top, Inches(7.5), Inches(0.9), WHITE)
        textbox(s, Inches(0.7), top + Inches(0.1), Inches(7.0), Inches(0.32), a, size=18, bold=True, color=FOREST)
        textbox(s, Inches(0.7), top + Inches(0.45), Inches(7.0), Inches(0.4), b, size=17, color=MUTED)
    round_rect(s, Inches(8.2), Inches(1.55), Inches(4.7), Inches(4.9), FOREST)
    textbox(s, Inches(8.5), Inches(1.8), Inches(4.2), Inches(0.4), "THE GAP IN TANZANIA", size=18, bold=True, color=LAGOON_BRIGHT)
    bullets(
        s,
        Inches(8.5),
        Inches(2.4),
        Inches(4.2),
        Inches(3.7),
        [
            "Few tools join live SME transactions with ML scoring.",
            "Most published benchmarks use developed-market data.",
            "Local supply-chain signals are still under-used.",
            "Ushirika turns that gap into a working prototype.",
        ],
        size=17,
        color=WHITE,
    )
    footer(s, 5)

    # —— 6 Method ——
    s = prs.slides.add_slide(blank)
    bg(s)
    section_header(s, "EVALUATION METHOD", "How we trained and tested the models")
    round_rect(s, Inches(0.45), Inches(1.55), Inches(6.15), Inches(5.0), WHITE)
    textbox(s, Inches(0.75), Inches(1.75), Inches(5.6), Inches(0.4), "TRAINING PROTOCOL", size=18, bold=True, color=LAGOON)
    bullets(
        s,
        Inches(0.75),
        Inches(2.3),
        Inches(5.6),
        Inches(4.0),
        [
            "Features from payment behaviour, volume, delays and partners.",
            "80/20 stratified train–test split (seed 42).",
            "Models fit on training data only; GridSearchCV uses ROC-AUC.",
            "Hold-out metrics: Accuracy, Precision, Recall, F1, ROC-AUC.",
            "An SME is scored only after at least twelve transactions.",
            "Rare large deals are down-weighted so they do not inflate loan size.",
        ],
        size=17,
    )
    round_rect(s, Inches(6.85), Inches(1.55), Inches(6.0), Inches(5.0), WHITE)
    textbox(s, Inches(7.15), Inches(1.75), Inches(5.5), Inches(0.4), "TWO MODELS COMPARED", size=18, bold=True, color=LAGOON)
    bullets(
        s,
        Inches(7.15),
        Inches(2.3),
        Inches(5.5),
        Inches(4.0),
        [
            "Baseline: Logistic Regression with scaling.",
            "Primary model: Random Forest Classifier.",
            "Selection metric: hold-out ROC-AUC.",
            "Score range about 300–850.",
            "Risk bands: Low ≥650 · Medium 500–649 · High <500.",
            "Users see plain signals such as on-time payments and late days.",
        ],
        size=17,
    )
    footer(s, 6)

    # —— 7 Findings (with evaluation proof) ——
    s = prs.slides.add_slide(blank)
    bg(s)
    section_header(
        s,
        "FINDINGS  ·  HOLD-OUT PROOF",
        "Random Forest beat Logistic Regression on unseen test data",
        "Hold-out evaluation  ·  80/20 stratified split  ·  test n = 280 SMEs",
    )

    # Comparison table
    round_rect(s, Inches(0.4), Inches(1.5), Inches(8.4), Inches(3.4), WHITE)
    textbox(s, Inches(0.6), Inches(1.6), Inches(8.0), Inches(0.35), "Hold-out metrics (creditworthy = positive class)", size=16, bold=True, color=FOREST)

    headers = ["Model", "Acc.", "Prec.", "Recall", "F1", "AUC"]
    rf_row = ["Random Forest", "88.6%", "93.0%", "85.7%", "89.2%", "95.8%"]
    lr_row = ["Logistic Reg.", "77.9%", "75.8%", "87.7%", "81.3%", "87.5%"]
    col_w = [1.85, 1.15, 1.2, 1.2, 1.1, 1.2]
    x0 = 0.65
    rect(s, Inches(0.55), Inches(2.05), Inches(8.05), Inches(0.48), FOREST)
    x = x0
    for h, w in zip(headers, col_w):
        textbox(s, Inches(x), Inches(2.12), Inches(w), Inches(0.35), h, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += w
    rect(s, Inches(0.55), Inches(2.53), Inches(8.05), Inches(0.55), SOFT)
    x = x0
    for val, w in zip(rf_row, col_w):
        textbox(s, Inches(x), Inches(2.62), Inches(w), Inches(0.4), val, size=16, bold=True, color=FOREST, align=PP_ALIGN.CENTER)
        x += w
    x = x0
    for val, w in zip(lr_row, col_w):
        textbox(s, Inches(x), Inches(3.2), Inches(w), Inches(0.4), val, size=16, color=MUTED, align=PP_ALIGN.CENTER)
        x += w
    textbox(
        s,
        Inches(0.65),
        Inches(3.75),
        Inches(7.9),
        Inches(0.95),
        "Proof of choice: RF ROC-AUC is +8.3 pts vs LR (95.8% vs 87.5%).\n"
        "RF accuracy +10.7 pts and precision +17.2 pts — fewer false “creditworthy” labels.",
        size=16,
        color=INK,
    )

    # Confusion matrix
    round_rect(s, Inches(9.0), Inches(1.5), Inches(3.85), Inches(3.4), FOREST)
    textbox(s, Inches(9.2), Inches(1.62), Inches(3.5), Inches(0.35), "RF confusion matrix", size=16, bold=True, color=LAGOON_BRIGHT)
    textbox(s, Inches(9.2), Inches(2.0), Inches(3.5), Inches(0.3), "Test set: 280 SMEs", size=15, color=SOFT)
    cm = [
        ("TN 116", "Correct higher-risk"),
        ("FP 10", "Wrongly called safe"),
        ("FN 22", "Missed creditworthy"),
        ("TP 132", "Correct creditworthy"),
    ]
    for i, (lab, tip) in enumerate(cm):
        top = Inches(2.45 + i * 0.55)
        textbox(s, Inches(9.2), top, Inches(3.5), Inches(0.28), lab, size=17, bold=True, color=WHITE)
        textbox(s, Inches(9.2), top + Inches(0.26), Inches(3.5), Inches(0.28), tip, size=14, color=RGBColor(0x9B, 0xC4, 0xBA))

    # Feature importance
    round_rect(s, Inches(0.4), Inches(5.05), Inches(12.5), Inches(1.7), WHITE)
    textbox(s, Inches(0.65), Inches(5.15), Inches(12.0), Inches(0.32), "What drove RF decisions (feature importance)", size=16, bold=True, color=FOREST)
    feats = [
        ("On-time rate", "19.8%"),
        ("Default rate", "18.7%"),
        ("Order-type mix", "8.3%"),
        ("Account age", "5.5%"),
        ("Buyer share", "4.6%"),
        ("Tx frequency", "4.4%"),
    ]
    for i, (name, pct) in enumerate(feats):
        left = Inches(0.6 + i * 2.1)
        textbox(s, left, Inches(5.55), Inches(2.0), Inches(0.35), pct, size=20, bold=True, color=LAGOON, align=PP_ALIGN.CENTER)
        textbox(s, left, Inches(6.0), Inches(2.0), Inches(0.45), name, size=15, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, 7)

    # —— 8 Results + limits / next steps (TRA + NIDA) ——
    s = prs.slides.add_slide(blank)
    bg(s)
    section_header(s, "RESULTS & NEXT STEPS", "Objectives met — and what must come next")

    rows = [
        ("1", "Features from data", "Transactions become behavioural predictors used for scoring."),
        ("2", "RF vs classical", "RF beats LR on ROC-AUC (95.8% vs 87.5%)."),
        ("3", "Standard evaluation", "Accuracy, Precision, Recall, F1 and ROC-AUC guide model choice."),
    ]
    for i, (num, title, body) in enumerate(rows):
        left = Inches(0.45 + i * 4.2)
        round_rect(s, left, Inches(1.55), Inches(4.0), Inches(1.9), WHITE)
        rect(s, left, Inches(1.55), Inches(4.0), Inches(0.48), FOREST if i != 1 else LAGOON)
        textbox(s, left + Inches(0.15), Inches(1.62), Inches(2.7), Inches(0.35), f"{num}  {title}", size=15, bold=True, color=WHITE)
        met_badge(s, left + Inches(2.9), Inches(1.62))
        textbox(s, left + Inches(0.2), Inches(2.2), Inches(3.6), Inches(1.05), body, size=16, color=MUTED)

    round_rect(s, Inches(0.45), Inches(3.65), Inches(12.4), Inches(2.95), FOREST)
    textbox(s, Inches(0.75), Inches(3.78), Inches(11.9), Inches(0.35), "LIMITS & NEXT STEPS", size=17, bold=True, color=LAGOON_BRIGHT)
    textbox(
        s,
        Inches(0.75),
        Inches(4.2),
        Inches(11.9),
        Inches(0.55),
        "The portal is live, but it is still a prototype — not a bank production system. "
        "Next: lender pilot on real SME ledgers, plus official verification links.",
        size=16,
        color=WHITE,
    )
    auth = [
        ("TRA — TIN", "Verify TIN numbers with TRA so registered businesses can be confirmed."),
        ("TRA — EFD", "Check EFD receipts so reported sales match fiscal records."),
        ("NIDA", "Verify national IDs with NIDA beyond format and date-of-birth checks."),
    ]
    for i, (title, body) in enumerate(auth):
        left = Inches(0.75 + i * 4.0)
        round_rect(s, left, Inches(4.9), Inches(3.8), Inches(1.45), RGBColor(0x0A, 0x35, 0x2A))
        textbox(s, left + Inches(0.2), Inches(5.0), Inches(3.4), Inches(0.3), title, size=16, bold=True, color=LAGOON_BRIGHT)
        textbox(s, left + Inches(0.2), Inches(5.4), Inches(3.4), Inches(0.8), body, size=15, color=MIST)
    footer(s, 8)

    # —— 9 Conclusion ——
    s = prs.slides.add_slide(blank)
    bg(s, FOREST_DEEP)
    rect(s, Inches(0), Inches(0), Inches(0.4), SLIDE_H, LAGOON_BRIGHT)
    textbox(s, Inches(0.85), Inches(0.5), Inches(11.5), Inches(0.35), "CONCLUSION", size=16, bold=True, color=LAGOON_BRIGHT)
    textbox(
        s,
        Inches(0.85),
        Inches(1.0),
        Inches(11.5),
        Inches(0.7),
        "The general objective was met",
        size=34,
        bold=True,
        color=WHITE,
        font="Georgia",
    )
    textbox(
        s,
        Inches(0.85),
        Inches(1.85),
        Inches(11.5),
        Inches(1.1),
        "Ushirika turns supply-chain transactions into clear credit scores. "
        "We compared Random Forest with Logistic Regression on hold-out data, "
        "and shipped a live portal for SMEs, lenders and admins.",
        size=19,
        color=MIST,
    )
    for i, (t, d) in enumerate(
        [
            ("Highest priority", "Pilot with partner lenders on real SME ledgers."),
            ("Authority links", "TRA (TIN & EFD) and NIDA verification for stronger trust."),
            ("Before production", "Managed database, hosting hardening, and ongoing model review."),
        ]
    ):
        top = Inches(3.25 + i * 0.85)
        round_rect(s, Inches(0.85), top, Inches(11.6), Inches(0.7), RGBColor(0x0A, 0x35, 0x2A))
        textbox(s, Inches(1.1), top + Inches(0.16), Inches(3.2), Inches(0.4), t, size=18, bold=True, color=LAGOON_BRIGHT)
        textbox(s, Inches(4.5), top + Inches(0.16), Inches(7.7), Inches(0.4), d, size=18, color=WHITE)
    textbox(s, Inches(0.85), Inches(6.45), Inches(10), Inches(0.4), "Asanteni  ·  Questions and discussion", size=20, bold=True, color=LAGOON_BRIGHT)
    textbox(s, Inches(11.2), Inches(6.45), Inches(1.7), Inches(0.4), "9  /  9", size=16, color=RGBColor(0xB8, 0xD0, 0xC8), align=PP_ALIGN.RIGHT)

    for path in (OUT, OUT_SEM2):
        try:
            path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(path))
            print(f"Saved: {path}")
        except PermissionError:
            if path == OUT:
                prs.save(str(OUT_ALT))
                print(f"Original PPT locked; saved: {OUT_ALT}")
            else:
                print(f"Could not write {path} (file may be open).")


if __name__ == "__main__":
    build()
