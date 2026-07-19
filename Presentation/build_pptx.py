"""
Build a professional Ushirika project presentation (PowerPoint).
Brand: lagoon teal + deep forest (matches the live portal).
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent
SHOTS = ROOT / "screenshots"
OUT = ROOT / "Ushirika_Project_Presentation.pptx"

# Brand palette (from Frontend styles)
FOREST = RGBColor(0x0B, 0x3D, 0x2E)
FOREST_DEEP = RGBColor(0x06, 0x28, 0x20)
LAGOON = RGBColor(0x1A, 0x7A, 0x6D)
LAGOON_BRIGHT = RGBColor(0x23, 0x96, 0x88)
MIST = RGBColor(0xE6, 0xEE, 0xF0)
PAPER = RGBColor(0xF4, 0xF8, 0xF9)
INK = RGBColor(0x0F, 0x28, 0x30)
INK_MUTED = RGBColor(0x4A, 0x5E, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x1A, 0x6B, 0x45)
WARN = RGBColor(0x8A, 0x5A, 0x00)
DANGER = RGBColor(0x9B, 0x2C, 0x2C)
LINE = RGBColor(0xC5, 0xD2, 0xD5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, text, size=18, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    # Support multi-line text
    lines = text.split("\n") if text is not None else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        set_run(run, line, size, bold, color, font)
    return box


def fill_shape(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_shape(shape, color)
    return shape


def rounded(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(shape, color)
    # softer corners
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def accent_bar(slide, left, top, width=Inches(0.12), height=Inches(0.55), color=LAGOON_BRIGHT):
    return rect(slide, left, top, width, height, color)


def footer(slide, page_num, total=12):
    rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), FOREST_DEEP)
    add_textbox(
        slide, Inches(0.45), Inches(7.18), Inches(8), Inches(0.28),
        "Ushirika  ·  Tanzania SME Credit Risk Platform",
        size=11, color=RGBColor(0xB8, 0xD0, 0xC8), font="Calibri",
    )
    add_textbox(
        slide, Inches(11.2), Inches(7.18), Inches(1.7), Inches(0.28),
        f"{page_num}  /  {total}",
        size=11, color=RGBColor(0xB8, 0xD0, 0xC8), align=PP_ALIGN.RIGHT, font="Calibri",
    )


def slide_bg(slide, color=PAPER):
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, color)


def title_block(slide, eyebrow, title, subtitle=None):
    add_textbox(slide, Inches(0.55), Inches(0.28), Inches(10), Inches(0.3), eyebrow, size=12, bold=True, color=LAGOON, font="Calibri")
    accent_bar(slide, Inches(0.55), Inches(0.62), Inches(0.1), Inches(0.5))
    add_textbox(slide, Inches(0.8), Inches(0.55), Inches(11.5), Inches(0.55), title, size=28, bold=True, color=FOREST, font="Calibri")
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(1.15), Inches(12), Inches(0.4), subtitle, size=14, color=INK_MUTED, font="Calibri")


def bullet_box(slide, left, top, width, height, items, size=15, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_before = Pt(5)
        p.space_after = Pt(5)
        run = p.add_run()
        set_run(run, f"•  {item}", size=size, color=color, font="Calibri")
    return box


def metric_card(slide, left, top, width, height, value, label, accent=LAGOON):
    rounded(slide, left, top, width, height, WHITE)
    rect(slide, left, top, Inches(0.1), height, accent)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.22), width - Inches(0.35), Inches(0.45), value, size=26, bold=True, color=FOREST, font="Calibri")
    add_textbox(slide, left + Inches(0.25), top + Inches(0.7), width - Inches(0.35), Inches(0.45), label, size=12, color=INK_MUTED, font="Calibri")


def add_image_fit(slide, path: Path, left, top, width, height):
    if not path.exists():
        rounded(slide, left, top, width, height, MIST)
        add_textbox(slide, left + Inches(0.2), top + height / 2 - Inches(0.2), width - Inches(0.4), Inches(0.4),
                    f"[Screenshot missing: {path.name}]", size=12, color=INK_MUTED, align=PP_ALIGN.CENTER)
        return None
    # Frame
    rounded(slide, left - Inches(0.06), top - Inches(0.06), width + Inches(0.12), height + Inches(0.12), LINE)
    pic = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    return pic


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = 12

    # ─── 1. Title ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    slide_bg(s, FOREST_DEEP)
    # decorative bands
    rect(s, Inches(0), Inches(0), Inches(0.35), SLIDE_H, LAGOON)
    rect(s, Inches(0), Inches(6.55), SLIDE_W, Inches(0.95), FOREST)
    add_textbox(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.35), "PROJECT PRESENTATION", size=13, bold=True, color=LAGOON_BRIGHT)
    add_textbox(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.1), "Ushirika", size=54, bold=True, color=WHITE, font="Georgia")
    add_textbox(
        s, Inches(0.8), Inches(3.15), Inches(11.5), Inches(0.9),
        "ML-Powered SME Credit Risk Assessment\nfor Inclusive Ecosystem Banking in Tanzania",
        size=22, color=MIST, font="Calibri",
    )
    add_textbox(
        s, Inches(0.8), Inches(4.5), Inches(11), Inches(0.7),
        "SME  ·  Lender  ·  Admin Portal  |  FastAPI  ·  Random Forest  ·  Vite",
        size=14, color=RGBColor(0x9B, 0xC4, 0xBA),
    )
    add_textbox(s, Inches(0.8), Inches(6.7), Inches(10), Inches(0.35), "Live: ushirika-sme-portal.vercel.app   ·   API: ushirika-api.onrender.com", size=12, color=RGBColor(0xB8, 0xD0, 0xC8))
    add_textbox(s, Inches(11.0), Inches(6.7), Inches(1.9), Inches(0.35), "1  /  12", size=12, color=RGBColor(0xB8, 0xD0, 0xC8), align=PP_ALIGN.RIGHT)

    # ─── 2. Agenda ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    slide_bg(s)
    title_block(s, "OVERVIEW", "Agenda", "What we will cover today")
    items = [
        ("01", "Problem & Opportunity", "SME credit access gaps in Tanzania"),
        ("02", "Solution Overview", "Ushirika platform at a glance"),
        ("03", "System Architecture", "Frontend, API, ML, deployment"),
        ("04", "User Roles & Journeys", "SME, Lender, Admin"),
        ("05", "Live Product Walkthrough", "Real screenshots from the portal"),
        ("06", "ML Credit Scoring", "Features, models, risk bands"),
        ("07", "Security & Impact", "Trust, privacy, outcomes"),
        ("08", "Demo & Next Steps", "Links, credentials, roadmap"),
    ]
    for i, (num, title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        left = Inches(0.55 + col * 6.3)
        top = Inches(1.75 + row * 1.2)
        rounded(s, left, top, Inches(5.95), Inches(1.0), WHITE)
        rect(s, left, top, Inches(0.12), Inches(1.0), LAGOON if i % 2 == 0 else FOREST)
        add_textbox(s, left + Inches(0.35), top + Inches(0.18), Inches(1), Inches(0.35), num, size=18, bold=True, color=LAGOON)
        add_textbox(s, left + Inches(1.2), top + Inches(0.15), Inches(4.4), Inches(0.35), title, size=16, bold=True, color=FOREST)
        add_textbox(s, left + Inches(1.2), top + Inches(0.5), Inches(4.4), Inches(0.35), desc, size=12, color=INK_MUTED)
    footer(s, 2, total)

    # ─── 3. Problem ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    slide_bg(s)
    title_block(s, "CONTEXT", "The Challenge", "Why traditional SME lending falls short")
    cards = [
        ("Limited credit history", "Many SMEs lack formal banking trails, so lenders struggle to assess risk fairly."),
        ("Manual underwriting", "Paper-based review is slow, inconsistent, and hard to scale across portfolios."),
        ("Supply-chain data unused", "Transaction patterns with buyers/suppliers contain strong predictive signals."),
        ("Inclusive finance gap", "Without data-driven scoring, viable businesses remain under-financed."),
    ]
    for i, (t, d) in enumerate(cards):
        left = Inches(0.55 + (i % 2) * 6.3)
        top = Inches(1.7 + (i // 2) * 2.35)
        rounded(s, left, top, Inches(6.0), Inches(2.1), WHITE)
        rect(s, left, top, Inches(6.0), Inches(0.12), LAGOON_BRIGHT)
        add_textbox(s, left + Inches(0.35), top + Inches(0.4), Inches(5.3), Inches(0.4), t, size=18, bold=True, color=FOREST)
        add_textbox(s, left + Inches(0.35), top + Inches(0.95), Inches(5.3), Inches(0.9), d, size=14, color=INK_MUTED)
    footer(s, 3, total)

    # ─── 4. Solution ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    slide_bg(s)
    title_block(s, "SOLUTION", "Introducing Ushirika", "Automated, data-driven ecosystem banking for SME credit risk")
    pillars = [
        ("For SMEs", "Record or import supply-chain transactions, view dashboards, and receive ML credit scores with eligible financing."),
        ("For Lenders", "Review SME portfolios by NIDA, inspect risk signals, and download statements for credit decisions."),
        ("For Admins", "Provision lenders, SMEs, and sub-admins; manage accounts; keep the ecosystem governed."),
    ]
    for i, (t, d) in enumerate(pillars):
        left = Inches(0.5 + i * 4.2)
        rounded(s, left, Inches(1.75), Inches(3.95), Inches(3.6), WHITE)
        rect(s, left, Inches(1.75), Inches(3.95), Inches(0.7), FOREST if i != 1 else LAGOON)
        add_textbox(s, left + Inches(0.25), Inches(1.9), Inches(3.45), Inches(0.45), t, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.3), Inches(2.75), Inches(3.35), Inches(2.3), d, size=14, color=INK)
    add_textbox(s, Inches(0.55), Inches(5.55), Inches(12), Inches(0.7),
                "Auth: NIDA / Membership ID + 4-digit PIN   ·   Financing capped at 75% of transaction volume   ·   Min. 5 transactions to score",
                size=13, color=INK_MUTED)
    footer(s, 4, total)

    # ─── 5. Architecture ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    slide_bg(s)
    title_block(s, "TECHNICAL DESIGN", "System Architecture", "End-to-end stack from browser to ML inference")
    layers = [
        ("Frontend", "Vite + Vanilla JS\nManrope / Newsreader UI\nRole-based dashboards", LAGOON),
        ("API Layer", "FastAPI + JWT\nSQLAlchemy ORM\nRole guards & CORS", FOREST),
        ("Data & ML", "SQLite / Postgres-ready\nRandom Forest + LR\n13 engineered features", LAGOON_BRIGHT),
        ("Cloud", "Vercel (portal)\nRender (API)\nAuto API switching", FOREST_DEEP),
    ]
    for i, (t, d, c) in enumerate(layers):
        left = Inches(0.5 + i * 3.2)
        rounded(s, left, Inches(1.8), Inches(3.0), Inches(3.5), WHITE)
        rect(s, left, Inches(1.8), Inches(3.0), Inches(0.65), c)
        add_textbox(s, left + Inches(0.15), Inches(1.92), Inches(2.7), Inches(0.4), t, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.25), Inches(2.7), Inches(2.5), Inches(2.3), d, size=13, color=INK, align=PP_ALIGN.CENTER)
        if i < 3:
            add_textbox(s, left + Inches(2.85), Inches(3.3), Inches(0.4), Inches(0.4), "→", size=22, bold=True, color=LAGOON, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.55), Inches(5.55), Inches(12), Inches(0.5),
                "Security: HMAC-SHA256 PII pseudonymization  ·  JWT expiry  ·  Soft-delete restore  ·  Auto logout on inactivity",
                size=13, color=INK_MUTED)
    footer(s, 5, total)

    # ─── 6. Roles ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    slide_bg(s)
    title_block(s, "USERS", "Roles & Capabilities", "One platform, three governed experiences")
    roles = [
        ("SME", "Self-register with NIDA + PIN", [
            "Create / edit / confirm transactions",
            "CSV import & e-statement export",
            "Credit score & financing estimate",
            "Profile & PIN management",
        ], LAGOON),
        ("Lender", "Created by Admin", [
            "Portfolio overview of SMEs",
            "Lookup by NIDA / membership",
            "View risk score & history",
            "Download SME statements",
        ], FOREST),
        ("Admin / Sub-Admin", "Bootstrap + delegated", [
            "CRUD for lenders & SMEs",
            "Create sub-admins",
            "Restore soft-deleted accounts",
            "Govern platform access",
        ], LAGOON_BRIGHT),
    ]
    for i, (name, how, bullets, accent) in enumerate(roles):
        left = Inches(0.45 + i * 4.25)
        rounded(s, left, Inches(1.7), Inches(4.05), Inches(4.7), WHITE)
        rect(s, left, Inches(1.7), Inches(4.05), Inches(1.05), accent)
        add_textbox(s, left + Inches(0.2), Inches(1.82), Inches(3.65), Inches(0.4), name, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.2), Inches(2.25), Inches(3.65), Inches(0.35), how, size=11, color=MIST, align=PP_ALIGN.CENTER)
        bullet_box(s, left + Inches(0.3), Inches(3.0), Inches(3.5), Inches(3.1), bullets, size=13)
    footer(s, 6, total)

    # ─── 7. Live walkthrough: Login + SME ────────────────────────
    s = prs.slides.add_slide(blank)
    slide_bg(s)
    title_block(s, "PRODUCT", "Live Walkthrough — Login & SME Portal", "Captured from the production Vercel deployment")
    add_image_fit(s, SHOTS / "01_login.png", Inches(0.45), Inches(1.65), Inches(6.0), Inches(3.75))
    add_image_fit(s, SHOTS / "02_sme_dashboard.png", Inches(6.85), Inches(1.65), Inches(6.0), Inches(3.75))
    add_textbox(s, Inches(0.45), Inches(5.55), Inches(6.0), Inches(0.35), "Secure NIDA / Membership + PIN login", size=12, bold=True, color=FOREST, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(6.85), Inches(5.55), Inches(6.0), Inches(0.35), "SME overview: score, financing & activity", size=12, bold=True, color=FOREST, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.55), Inches(6.05), Inches(12), Inches(0.45),
                "SMEs build a verified transaction history — the foundation for transparent, model-driven credit decisions.",
                size=13, color=INK_MUTED)
    footer(s, 7, total)

    # ─── 8. SME transactions + Lender ───────────────────────────
    s = prs.slides.add_slide(blank)
    slide_bg(s)
    title_block(s, "PRODUCT", "Transactions, Portfolio & Credit Insight", "Operational workflows for SMEs and lenders")
    add_image_fit(s, SHOTS / "03_sme_transactions.png", Inches(0.4), Inches(1.6), Inches(6.1), Inches(3.8))
    add_image_fit(s, SHOTS / "05_lender_portfolio.png", Inches(6.8), Inches(1.6), Inches(6.1), Inches(3.8))
    add_textbox(s, Inches(0.4), Inches(5.55), Inches(6.1), Inches(0.35), "SME transaction ledger & controls", size=12, bold=True, color=FOREST, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(6.8), Inches(5.55), Inches(6.1), Inches(0.35), "Lender portfolio & NIDA lookup", size=12, bold=True, color=FOREST, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.55), Inches(6.05), Inches(12), Inches(0.45),
                "CSV upload/export available  ·  Auto-rescore after create / update / delete / import when ≥ 5 transactions",
                size=13, color=INK_MUTED)
    footer(s, 8, total)

    # ─── 9. Admin + detail ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    slide_bg(s)
    title_block(s, "PRODUCT", "Admin Governance & Lender Deep-Dive", "Account lifecycle and SME risk detail views")
    img_a = SHOTS / "07_admin_accounts.png"
    img_b = SHOTS / "06_lender_sme_detail.png"
    if not img_b.exists():
        img_b = SHOTS / "08_admin_create.png"
    add_image_fit(s, img_a, Inches(0.4), Inches(1.6), Inches(6.1), Inches(3.8))
    add_image_fit(s, img_b, Inches(6.8), Inches(1.6), Inches(6.1), Inches(3.8))
    add_textbox(s, Inches(0.4), Inches(5.55), Inches(6.1), Inches(0.35), "Admin account management", size=12, bold=True, color=FOREST, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(6.8), Inches(5.55), Inches(6.1), Inches(0.35), "Lender SME detail / create flow", size=12, bold=True, color=FOREST, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.55), Inches(6.05), Inches(12), Inches(0.45),
                "Admins create lenders, SMEs, and sub-admins; soft-deleted accounts can be restored without data loss.",
                size=13, color=INK_MUTED)
    footer(s, 9, total)

    # ─── 10. ML Scoring ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    slide_bg(s)
    title_block(s, "INTELLIGENCE", "ML Credit Scoring Engine", "Random Forest vs Logistic Regression — honest hold-out evaluation")

    # Metrics
    metric_card(s, Inches(0.5), Inches(1.65), Inches(3.0), Inches(1.35), "88.5%", "RF ROC-AUC (primary)", SUCCESS)
    metric_card(s, Inches(3.7), Inches(1.65), Inches(3.0), Inches(1.35), "84.2%", "RF Accuracy", LAGOON)
    metric_card(s, Inches(6.9), Inches(1.65), Inches(3.0), Inches(1.35), "75.4%", "LR ROC-AUC (baseline)", WARN)
    metric_card(s, Inches(10.1), Inches(1.65), Inches(2.7), Inches(1.35), "RF wins", "Outperforms baseline", FOREST)

    rounded(s, Inches(0.5), Inches(3.25), Inches(6.1), Inches(3.2), WHITE)
    add_textbox(s, Inches(0.75), Inches(3.4), Inches(5.6), Inches(0.35), "13 engineered features", size=16, bold=True, color=FOREST)
    bullet_box(s, Inches(0.75), Inches(3.85), Inches(5.6), Inches(2.4), [
        "Payment consistency, delays, on-time rate",
        "Turnover (TZS), frequency, volume trend",
        "Default & compliance rates, completion",
        "Counterparty diversity, account age",
        "Avg. transaction interval (higher → riskier)",
    ], size=13)

    rounded(s, Inches(6.9), Inches(3.25), Inches(5.9), Inches(3.2), WHITE)
    add_textbox(s, Inches(7.15), Inches(3.4), Inches(5.4), Inches(0.35), "Score mapping & risk bands", size=16, bold=True, color=FOREST)
    bullet_box(s, Inches(7.15), Inches(3.85), Inches(5.4), Inches(2.4), [
        "Probability → conservative score (~300–680)",
        "Low ≥ 580  ·  Medium 480–579  ·  High < 480",
        "Eligible financing capped at 75% of volume",
        "Min / Max band: 500k – 50M TZS",
        "Requires ≥ 5 transactions before scoring",
    ], size=13)
    footer(s, 10, total)

    # ─── 11. Security & Impact ──────────────────────────────────
    s = prs.slides.add_slide(blank)
    slide_bg(s)
    title_block(s, "TRUST & VALUE", "Security, Privacy & Impact", "Built for responsible deployment in financial contexts")
    left_items = [
        ("Privacy by design", "Sensitive fields hashed with HMAC-SHA256 — no plain PII in analytics stores."),
        ("Access control", "JWT auth, role guards, PIN change, inactivity auto-logout."),
        ("Operational integrity", "Confirm / edit / delete transactions with audit-friendly flows."),
    ]
    right_items = [
        ("Faster decisions", "Lenders see scored portfolios instead of starting from blank paper files."),
        ("Fairer access", "Supply-chain behavior unlocks credit for SMEs without long bank histories."),
        ("Scalable ops", "CSV bulk import + cloud hosting (Vercel + Render) for demo-to-pilot readiness."),
    ]
    for i, (t, d) in enumerate(left_items):
        top = Inches(1.7 + i * 1.55)
        rounded(s, Inches(0.5), top, Inches(6.0), Inches(1.4), WHITE)
        rect(s, Inches(0.5), top, Inches(0.12), Inches(1.4), LAGOON)
        add_textbox(s, Inches(0.85), top + Inches(0.25), Inches(5.4), Inches(0.35), t, size=15, bold=True, color=FOREST)
        add_textbox(s, Inches(0.85), top + Inches(0.65), Inches(5.4), Inches(0.55), d, size=12, color=INK_MUTED)
    for i, (t, d) in enumerate(right_items):
        top = Inches(1.7 + i * 1.55)
        rounded(s, Inches(6.8), top, Inches(6.0), Inches(1.4), WHITE)
        rect(s, Inches(6.8), top, Inches(0.12), Inches(1.4), FOREST)
        add_textbox(s, Inches(7.15), top + Inches(0.25), Inches(5.4), Inches(0.35), t, size=15, bold=True, color=FOREST)
        add_textbox(s, Inches(7.15), top + Inches(0.65), Inches(5.4), Inches(0.55), d, size=12, color=INK_MUTED)
    footer(s, 11, total)

    # ─── 12. Closing / Demo ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    slide_bg(s, FOREST_DEEP)
    rect(s, Inches(0), Inches(0), Inches(0.35), SLIDE_H, LAGOON_BRIGHT)
    add_textbox(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.35), "THANK YOU", size=13, bold=True, color=LAGOON_BRIGHT)
    add_textbox(s, Inches(0.8), Inches(1.05), Inches(11.5), Inches(0.7), "Questions & Live Demo", size=36, bold=True, color=WHITE, font="Georgia")
    add_textbox(s, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.45),
                "Explore the live system — then let’s discuss the roadmap.", size=16, color=MIST)

    # Links card
    rounded(s, Inches(0.8), Inches(2.55), Inches(7.2), Inches(2.6), FOREST)
    add_textbox(s, Inches(1.05), Inches(2.75), Inches(6.7), Inches(0.35), "Live links", size=14, bold=True, color=LAGOON_BRIGHT)
    add_textbox(s, Inches(1.05), Inches(3.2), Inches(6.7), Inches(1.6),
                "Portal:  https://ushirika-sme-portal.vercel.app\n"
                "API:       https://ushirika-api.onrender.com\n"
                "Docs:     /api/docs on the API host\n"
                "GitHub:  github.com/Dmmkuliku/ushirika-sme-credit",
                size=14, color=WHITE)

    # Creds
    rounded(s, Inches(8.3), Inches(2.55), Inches(4.4), Inches(2.6), FOREST)
    add_textbox(s, Inches(8.55), Inches(2.75), Inches(4.0), Inches(0.35), "Demo PINs (all 1234)", size=14, bold=True, color=LAGOON_BRIGHT)
    add_textbox(s, Inches(8.55), Inches(3.25), Inches(4.0), Inches(1.6),
                "Admin   20031001121160000228\n"
                "Lender  EMP001\n"
                "SME       19900101123456789012",
                size=13, color=WHITE)

    # Roadmap
    add_textbox(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.35), "Next steps", size=14, bold=True, color=LAGOON_BRIGHT)
    add_textbox(s, Inches(0.8), Inches(5.85), Inches(11.5), Inches(0.7),
                "Managed Postgres  ·  Real NIDA verification  ·  Lender decision workflow  ·  Model monitoring dashboard  ·  Mobile-responsive polish",
                size=13, color=MIST)
    add_textbox(s, Inches(0.8), Inches(6.75), Inches(10), Inches(0.35), "Ushirika — Inclusive SME credit, powered by data.", size=12, color=RGBColor(0x9B, 0xC4, 0xBA))
    add_textbox(s, Inches(11.0), Inches(6.75), Inches(1.9), Inches(0.35), "12  /  12", size=12, color=RGBColor(0xB8, 0xD0, 0xC8), align=PP_ALIGN.RIGHT)

    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
